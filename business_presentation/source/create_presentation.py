from pathlib import Path
import math
import textwrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET_SVG = ROOT / "assets" / "svg"
ASSET_PNG = ROOT / "assets" / "png"
DOCS = ROOT / "docs"
OUT = ROOT

BLUE = "2563EB"
BLUE_DARK = "1D4ED8"
NAVY = "0F172A"
GRAY = "64748B"
LIGHT = "F8FAFC"
BORDER = "D9E3F0"
TEAL = "14B8A6"
GREEN = "16A34A"
RED = "DC2626"
WHITE = "FFFFFF"

COLORS = {
    "blue": BLUE,
    "blue_dark": BLUE_DARK,
    "navy": NAVY,
    "gray": GRAY,
    "light": LIGHT,
    "border": BORDER,
    "teal": TEAL,
    "green": GREEN,
    "red": RED,
    "white": WHITE,
}

PRODUCTS = [
    ("AR-CERT-PRO", "Generate hundreds of certificates in one click.", "certificate"),
    ("AR-WORKSHEET-PRO", "Generate printable worksheets instantly.", "worksheet"),
    ("AR-QUESTION-PRO", "Create question papers automatically.", "question"),
    ("AR-ATTENDANCE-PRO", "Attendance reports.", "attendance"),
    ("AR-IDCARD-PRO", "Student ID cards.", "idcard"),
    ("AR-RESULT-PRO", "Results and mark sheets.", "result"),
    ("AR-TIMETABLE-PRO", "Automatic timetable generation.", "timetable"),
]

PROBLEMS = [
    ("Certificate creation", "certificate"),
    ("Worksheet preparation", "worksheet"),
    ("Attendance", "attendance"),
    ("ID Cards", "idcard"),
    ("Result preparation", "result"),
    ("Question papers", "question"),
    ("Timetables", "timetable"),
]

BENEFITS = [
    ("Save Time", "time"),
    ("Reduce Errors", "shield"),
    ("Professional Documents", "document"),
    ("One Click", "click"),
    ("Reusable Templates", "template"),
    ("Printable PDFs", "pdf"),
]

AUDIENCES = [
    ("Coaching Institutes", "coaching"),
    ("Schools", "school"),
    ("Colleges", "college"),
    ("Training Institutes", "training"),
    ("HR Teams", "hr"),
    ("Certification Organizations", "cert_org"),
]


def ensure_dirs():
    for path in [ASSET_SVG, ASSET_PNG, DOCS]:
        path.mkdir(parents=True, exist_ok=True)


def rgb(hex_color):
    h = hex_color.strip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_run(run, size=18, bold=False, color=NAVY, font="Inter"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_text(slide, text, x, y, w, h, size=20, color=NAVY, bold=False, font="Inter",
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_round_rect(slide, x, y, w, h, fill=WHITE, line=BORDER, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    return shape


def add_button(slide, text, x, y, w, h):
    btn = add_round_rect(slide, x, y, w, h, fill=BLUE, line=BLUE)
    tf = btn.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_run(run, size=15, bold=True, color=WHITE, font="Poppins")
    return btn


def add_logo(slide, x, y, scale=1):
    group = slide.shapes.add_group_shape()
    group.left = Inches(x)
    group.top = Inches(y)
    group.width = Inches(3.6 * scale)
    group.height = Inches(0.62 * scale)
    mark = group.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0, 0, Inches(0.54 * scale), Inches(0.54 * scale))
    mark.fill.solid()
    mark.fill.fore_color.rgb = rgb(BLUE)
    mark.line.color.rgb = rgb(BLUE)
    txt = group.shapes.add_textbox(Inches(0.72 * scale), Inches(0.02 * scale), Inches(3.0 * scale), Inches(0.42 * scale))
    tf = txt.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "AaryaRushi"
    set_run(r, size=int(18 * scale), bold=True, color=NAVY, font="Poppins")
    accent = group.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.16 * scale), Inches(0.16 * scale), Inches(0.22 * scale), Inches(0.22 * scale))
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(WHITE)
    accent.line.color.rgb = rgb(WHITE)
    return group


def add_header(slide, title, subtitle=None):
    add_logo(slide, 0.6, 0.32, 0.72)
    add_text(slide, title, 0.7, 0.98, 8.4, 0.42, size=26, bold=True, font="Poppins")
    if subtitle:
        add_text(slide, subtitle, 0.72, 1.42, 8.8, 0.34, size=12, color=GRAY)


def icon_svg(name, stroke=BLUE, fill="none"):
    common = f'xmlns="http://www.w3.org/2000/svg" width="96" height="96" viewBox="0 0 96 96" fill="{fill}" stroke="#{stroke}" stroke-width="4" stroke-linecap="round" stroke-linejoin="round"'
    icons = {
        "certificate": '<path d="M24 18h48v44H24z"/><path d="M34 30h28M34 42h20"/><circle cx="62" cy="62" r="10"/><path d="m56 70-4 12 10-5 10 5-4-12"/>',
        "worksheet": '<path d="M26 16h44v64H26z"/><path d="M36 32h24M36 44h24M36 56h14"/><path d="M58 56h2"/>',
        "question": '<path d="M28 18h40v52H28z"/><path d="M40 34a8 8 0 1 1 13 6c-4 3-5 5-5 9"/><path d="M48 60h.1"/>',
        "attendance": '<path d="M24 22h48v52H24z"/><path d="M34 36h28M34 50h28"/><path d="m35 62 6 6 14-16"/>',
        "idcard": '<rect x="18" y="24" width="60" height="48" rx="6"/><circle cx="38" cy="44" r="8"/><path d="M52 38h14M52 50h14M30 62c4-7 12-7 16 0"/>',
        "result": '<path d="M24 18h48v60H24z"/><path d="M36 60V42M48 60V30M60 60V48"/><path d="M34 66h28"/>',
        "timetable": '<rect x="18" y="22" width="60" height="56" rx="6"/><path d="M18 38h60M34 22v56M50 38v40M66 38v40"/>',
        "time": '<circle cx="48" cy="48" r="30"/><path d="M48 30v20l14 8"/>',
        "shield": '<path d="M48 16 72 26v18c0 18-10 30-24 36-14-6-24-18-24-36V26z"/><path d="m36 48 8 8 18-20"/>',
        "document": '<path d="M28 16h28l14 14v50H28z"/><path d="M56 16v16h14M38 46h20M38 58h20"/>',
        "click": '<path d="M42 18v34"/><path d="m42 52 10-10M42 52 32 42"/><path d="M54 56h14l-8 20-6-12-10 8z"/>',
        "template": '<rect x="18" y="18" width="60" height="60" rx="6"/><path d="M18 36h60M36 36v42M48 50h18M48 62h12"/>',
        "pdf": '<path d="M28 16h28l14 14v50H28z"/><path d="M56 16v16h14"/><path d="M34 60h28"/><path d="M34 48h8c6 0 6-10 0-10h-8v22"/>',
        "school": '<path d="M16 44 48 26l32 18-32 18z"/><path d="M28 52v16c10 7 30 7 40 0V52"/>',
        "coaching": '<rect x="20" y="24" width="56" height="40" rx="4"/><path d="M32 76h32M48 64v12M32 38h32M32 50h20"/>',
        "college": '<path d="M18 40 48 20l30 20z"/><path d="M24 40v34M72 40v34M36 46v28M60 46v28M20 74h56"/>',
        "training": '<circle cx="36" cy="36" r="10"/><path d="M20 72c4-16 28-16 32 0"/><path d="M58 30h18M58 44h18M58 58h12"/>',
        "hr": '<circle cx="34" cy="34" r="9"/><circle cx="62" cy="34" r="9"/><path d="M18 72c4-15 28-15 32 0M46 72c4-15 28-15 32 0"/>',
        "cert_org": '<circle cx="48" cy="42" r="20"/><path d="m36 76 8-18M60 76l-8-18M38 42l7 7 14-16"/>',
    }
    return f"<svg {common}>{icons[name]}</svg>\n"


def save_svg_assets():
    names = sorted(set([i for _, i in PROBLEMS] + [i for _, _, i in PRODUCTS] + [i for _, i in BENEFITS] + [i for _, i in AUDIENCES]))
    for name in names:
        (ASSET_SVG / f"{name}.svg").write_text(icon_svg(name), encoding="utf-8")
    logo = f'''<svg xmlns="http://www.w3.org/2000/svg" width="720" height="150" viewBox="0 0 720 150">
<rect x="28" y="28" width="92" height="92" rx="22" fill="#{BLUE}"/>
<circle cx="74" cy="74" r="20" fill="#fff"/>
<text x="148" y="75" fill="#{NAVY}" font-family="Poppins, Inter, Arial" font-size="42" font-weight="700">AaryaRushi</text>
<text x="150" y="108" fill="#{GRAY}" font-family="Inter, Arial" font-size="22">Automation Labs</text>
</svg>
'''
    (ASSET_SVG / "aaryarushi_logo.svg").write_text(logo, encoding="utf-8")


def load_font(size, bold=False):
    candidates = [
        "C:/Windows/Fonts/Inter.ttf",
        "C:/Windows/Fonts/Poppins-Regular.ttf",
        "C:/Windows/Fonts/segoeui.ttf",
        "C:/Windows/Fonts/arial.ttf",
    ]
    bold_candidates = [
        "C:/Windows/Fonts/Inter-Bold.ttf",
        "C:/Windows/Fonts/Poppins-Bold.ttf",
        "C:/Windows/Fonts/segoeuib.ttf",
        "C:/Windows/Fonts/arialbd.ttf",
    ]
    for path in (bold_candidates if bold else candidates):
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def draw_png_icon(name, path, size=512, stroke=BLUE):
    im = Image.new("RGBA", (size, size), (255, 255, 255, 0))
    d = ImageDraw.Draw(im)
    s = size / 96
    col = tuple(int(stroke[i:i + 2], 16) for i in (0, 2, 4)) + (255,)
    lw = int(4 * s)

    def xy(vals):
        return tuple(int(v * s) for v in vals)

    # Minimal matching PNG icon set for presentation use.
    if name in {"certificate", "worksheet", "question", "attendance", "document", "pdf", "result"}:
        d.rounded_rectangle(xy((24, 14, 72, 80)), radius=int(4 * s), outline=col, width=lw)
        if name == "certificate":
            d.line(xy((34, 30, 62, 30)), fill=col, width=lw)
            d.line(xy((34, 42, 56, 42)), fill=col, width=lw)
            d.ellipse(xy((52, 54, 72, 74)), outline=col, width=lw)
        elif name == "question":
            d.arc(xy((38, 28, 58, 48)), 205, 510, fill=col, width=lw)
            d.line(xy((48, 48, 48, 56)), fill=col, width=lw)
            d.ellipse(xy((46, 64, 50, 68)), fill=col)
        elif name == "attendance":
            d.line(xy((34, 34, 62, 34)), fill=col, width=lw)
            d.line(xy((34, 48, 62, 48)), fill=col, width=lw)
            d.line(xy((35, 62, 42, 68, 58, 52)), fill=col, width=lw)
        elif name == "result":
            d.line(xy((36, 60, 36, 44)), fill=col, width=lw)
            d.line(xy((48, 60, 48, 30)), fill=col, width=lw)
            d.line(xy((60, 60, 60, 48)), fill=col, width=lw)
            d.line(xy((34, 66, 62, 66)), fill=col, width=lw)
        else:
            d.line(xy((34, 34, 62, 34)), fill=col, width=lw)
            d.line(xy((34, 48, 62, 48)), fill=col, width=lw)
            d.line(xy((34, 62, 54, 62)), fill=col, width=lw)
    elif name == "idcard":
        d.rounded_rectangle(xy((18, 24, 78, 72)), radius=int(6 * s), outline=col, width=lw)
        d.ellipse(xy((30, 36, 46, 52)), outline=col, width=lw)
        d.line(xy((54, 38, 68, 38)), fill=col, width=lw)
        d.line(xy((54, 50, 68, 50)), fill=col, width=lw)
    elif name == "timetable":
        d.rounded_rectangle(xy((18, 22, 78, 78)), radius=int(6 * s), outline=col, width=lw)
        d.line(xy((18, 38, 78, 38)), fill=col, width=lw)
        d.line(xy((34, 22, 34, 78)), fill=col, width=lw)
        d.line(xy((50, 38, 50, 78)), fill=col, width=lw)
        d.line(xy((66, 38, 66, 78)), fill=col, width=lw)
    elif name == "time":
        d.ellipse(xy((18, 18, 78, 78)), outline=col, width=lw)
        d.line(xy((48, 30, 48, 50, 62, 58)), fill=col, width=lw)
    elif name == "shield":
        d.polygon([xy((48, 16))[0:2], xy((72, 26))[0:2], xy((70, 52))[0:2], xy((48, 80))[0:2], xy((26, 52))[0:2], xy((24, 26))[0:2]], outline=col)
        d.line(xy((36, 48, 44, 56, 62, 36)), fill=col, width=lw)
    elif name == "click":
        d.line(xy((42, 18, 42, 52, 58, 38)), fill=col, width=lw)
        d.polygon([xy((54, 56))[0:2], xy((72, 56))[0:2], xy((62, 78))[0:2]], outline=col)
    elif name == "template":
        d.rounded_rectangle(xy((18, 18, 78, 78)), radius=int(6 * s), outline=col, width=lw)
        d.line(xy((18, 36, 78, 36)), fill=col, width=lw)
        d.line(xy((36, 36, 36, 78)), fill=col, width=lw)
    else:
        d.rounded_rectangle(xy((20, 24, 76, 70)), radius=int(6 * s), outline=col, width=lw)
        d.line(xy((32, 42, 64, 42)), fill=col, width=lw)
        d.line(xy((32, 54, 56, 54)), fill=col, width=lw)
    im.save(path)


def save_png_assets():
    names = sorted(set([i for _, i in PROBLEMS] + [i for _, _, i in PRODUCTS] + [i for _, i in BENEFITS] + [i for _, i in AUDIENCES]))
    for name in names:
        draw_png_icon(name, ASSET_PNG / f"{name}.png")

    logo = Image.new("RGBA", (1200, 280), (255, 255, 255, 0))
    d = ImageDraw.Draw(logo)
    d.rounded_rectangle((40, 52, 196, 208), radius=34, fill=tuple(int(BLUE[i:i + 2], 16) for i in (0, 2, 4)) + (255,))
    d.ellipse((88, 100, 148, 160), fill=(255, 255, 255, 255))
    d.text((240, 66), "AaryaRushi", font=load_font(72, True), fill=tuple(int(NAVY[i:i + 2], 16) for i in (0, 2, 4)))
    d.text((244, 150), "Automation Labs", font=load_font(38), fill=tuple(int(GRAY[i:i + 2], 16) for i in (0, 2, 4)))
    logo.save(ASSET_PNG / "aaryarushi_logo.png")

    qr = Image.new("RGBA", (640, 640), (255, 255, 255, 255))
    d = ImageDraw.Draw(qr)
    d.rounded_rectangle((0, 0, 639, 639), radius=36, outline=tuple(int(BORDER[i:i + 2], 16) for i in (0, 2, 4)), width=8)
    cell = 32
    pattern = [
        "11111110011001111111",
        "10000010100101000001",
        "10111010111101011101",
        "10111010001001011101",
        "10111010101001011101",
        "10000010110101000001",
        "11111110101001111111",
        "00000000111100000000",
        "11010111100111101011",
        "00101100111000110100",
        "11100011101110100111",
        "01011001010001110010",
        "10111101110111011101",
        "00000000101100000000",
        "11111110100101111111",
        "10000010111101000001",
        "10111010011001011101",
        "10111010100101011101",
        "10000010111001000001",
        "11111110100101111111",
    ]
    ox = oy = 0
    dark = tuple(int(NAVY[i:i + 2], 16) for i in (0, 2, 4)) + (255,)
    for r, row in enumerate(pattern):
        for c, val in enumerate(row):
            if val == "1":
                d.rounded_rectangle((ox + c * cell + 6, oy + r * cell + 6, ox + (c + 1) * cell - 6, oy + (r + 1) * cell - 6), radius=4, fill=dark)
    d.text((190, 294), "QR", font=load_font(48, True), fill=tuple(int(BLUE[i:i + 2], 16) for i in (0, 2, 4)))
    d.text((142, 354), "placeholder", font=load_font(28), fill=tuple(int(GRAY[i:i + 2], 16) for i in (0, 2, 4)))
    qr.save(ASSET_PNG / "qr_placeholder.png")


def add_icon(slide, name, x, y, w=0.34):
    slide.shapes.add_picture(str(ASSET_PNG / f"{name}.png"), Inches(x), Inches(y), width=Inches(w))


def add_manual_work_illustration(slide, x, y, w, h):
    card = add_round_rect(slide, x, y, w, h, fill="FFFFFF", line=BORDER)
    add_text(slide, "Manual office work", x + 0.35, y + 0.28, w - 0.7, 0.28, size=13, bold=True)
    add_text(slide, "Repeated typing, formatting and checking", x + 0.35, y + 0.62, w - 0.7, 0.3, size=10, color=GRAY)
    desk = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.42), Inches(y + h - 0.78), Inches(w - 0.84), Inches(0.12))
    desk.fill.solid()
    desk.fill.fore_color.rgb = rgb(BORDER)
    desk.line.color.rgb = rgb(BORDER)
    for i, dx in enumerate([0.55, 1.75, 2.95]):
        paper = add_round_rect(slide, x + dx, y + 1.05 + (i % 2) * 0.25, 0.84, 1.05, fill=LIGHT, line=BORDER, radius=False)
        for j in range(3):
            ln = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + dx + 0.12), Inches(y + 1.23 + (i % 2) * 0.25 + j * 0.2), Inches(0.54), Inches(0.03))
            ln.fill.solid()
            ln.fill.fore_color.rgb = rgb(GRAY if j == 0 else BORDER)
            ln.line.color.rgb = ln.fill.fore_color.rgb
    head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + w - 1.35), Inches(y + 1.15), Inches(0.42), Inches(0.42))
    head.fill.solid()
    head.fill.fore_color.rgb = rgb(BLUE)
    head.line.color.rgb = rgb(BLUE)
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + w - 1.52), Inches(y + 1.62), Inches(0.76), Inches(0.72))
    body.fill.solid()
    body.fill.fore_color.rgb = rgb(TEAL)
    body.line.color.rgb = rgb(TEAL)
    return card


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)
    add_logo(slide, 0.78, 0.55, 1.0)
    add_text(slide, "Office Automation Solutions", 1.52, 1.22, 3.2, 0.28, size=12, color=GRAY)
    add_text(slide, "Save Hours of Manual Office Work Every Week", 0.9, 2.0, 6.9, 1.35, size=43, bold=True, font="Poppins")
    add_text(slide, "Smart Office Automation for Schools, Coaching Institutes and Training Centers", 0.94, 3.55, 6.6, 0.62, size=18, color=GRAY)
    add_button(slide, "Book Free 20-Min Audit", 0.94, 4.48, 2.55, 0.54)
    add_round_rect(slide, 8.25, 1.08, 3.75, 4.95, fill=LIGHT, line=BORDER)
    add_text(slide, "Automation Engine", 8.58, 1.44, 2.4, 0.32, size=18, bold=True, font="Poppins")
    add_text(slide, "Excel data to ready PDFs", 8.6, 1.82, 2.3, 0.25, size=11, color=GRAY)
    for i, label in enumerate(["Excel", "Word", "PDF"]):
        add_round_rect(slide, 8.72 + i * 1.0, 2.55, 0.78, 0.78, fill=WHITE, line=BORDER)
        add_text(slide, label, 8.74 + i * 1.0, 2.82, 0.74, 0.18, size=9, bold=True, align=PP_ALIGN.CENTER)
    add_round_rect(slide, 8.76, 4.05, 2.55, 0.32, fill="E0F2FE", line="E0F2FE")
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.76), Inches(4.05), Inches(1.98), Inches(0.32))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(BLUE)
    bar.line.color.rgb = rgb(BLUE)
    add_text(slide, "300 documents completed", 8.76, 4.55, 2.4, 0.28, size=11, color=GREEN, bold=True)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Current Problems", "Daily office work still consumes hours of staff time.")
    add_manual_work_illustration(slide, 8.1, 1.48, 4.15, 4.65)
    for idx, (label, icon) in enumerate(PROBLEMS):
        col = idx % 2
        row = idx // 2
        x = 0.82 + col * 3.3
        y = 2.0 + row * 0.86
        add_round_rect(slide, x, y, 2.9, 0.58, fill=WHITE, line=BORDER)
        add_text(slide, "X", x + 0.16, y + 0.17, 0.25, 0.2, size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)
        add_icon(slide, icon, x + 0.52, y + 0.12, 0.28)
        add_text(slide, label, x + 0.92, y + 0.18, 1.8, 0.2, size=12, bold=True)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_header(slide, "What We Automate", "Ready-made automation systems for common institute workflows.")
    for idx, (title, desc, icon) in enumerate(PRODUCTS):
        col = idx % 3
        row = idx // 3
        x = 0.72 + col * 4.08
        y = 1.95 + row * 1.48
        w = 3.55
        add_round_rect(slide, x, y, w, 1.12, fill=WHITE, line=BORDER)
        add_icon(slide, icon, x + 0.2, y + 0.23, 0.42)
        add_text(slide, title, x + 0.78, y + 0.25, 2.55, 0.24, size=12, bold=True, font="Poppins")
        add_text(slide, desc, x + 0.78, y + 0.57, 2.45, 0.38, size=10.5, color=GRAY)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_header(slide, "How It Works", "A simple office workflow: data in, documents out.")
    flow = [("Excel", "worksheet"), ("Automation Engine", "template"), ("Word", "document"), ("PDF", "pdf"), ("Ready", "shield")]
    for i, (label, icon) in enumerate(flow):
        x = 0.9 + i * 2.45
        add_round_rect(slide, x, 2.52, 1.55, 1.25, fill=WHITE, line=BORDER)
        add_icon(slide, icon, x + 0.52, 2.75, 0.42)
        add_text(slide, label, x + 0.15, 3.35, 1.25, 0.22, size=11, bold=True, align=PP_ALIGN.CENTER)
        if i < len(flow) - 1:
            add_text(slide, ">", x + 1.72, 2.95, 0.42, 0.42, size=25, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Clean templates keep every output professional, consistent and ready to share.", 2.2, 4.72, 8.9, 0.42, size=17, color=GRAY, align=PP_ALIGN.CENTER)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Example: 300 Student PDFs", "From spreadsheet rows to completed documents.")
    add_round_rect(slide, 0.82, 1.82, 3.4, 3.5, fill=WHITE, line=BORDER)
    add_text(slide, "Excel contains", 1.12, 2.2, 2.2, 0.26, size=14, color=GRAY)
    add_text(slide, "300 students", 1.12, 2.58, 2.2, 0.45, size=28, bold=True, font="Poppins")
    add_button(slide, "Generate", 1.12, 3.68, 1.5, 0.5)
    add_round_rect(slide, 5.0, 1.82, 3.4, 3.5, fill=WHITE, line=BORDER)
    add_text(slide, "Generating...", 5.32, 2.18, 2.2, 0.3, size=16, bold=True)
    add_round_rect(slide, 5.32, 2.84, 2.55, 0.28, fill="E0F2FE", line="E0F2FE")
    progress = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(5.32), Inches(2.84), Inches(2.15), Inches(0.28))
    progress.fill.solid()
    progress.fill.fore_color.rgb = rgb(BLUE)
    progress.line.color.rgb = rgb(BLUE)
    for i, num in enumerate(["25", "75", "150", "300"]):
        add_text(slide, num, 5.32 + i * 0.68, 3.42, 0.42, 0.22, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "Completed", 5.32, 4.2, 1.6, 0.3, size=15, color=GREEN, bold=True)
    add_round_rect(slide, 9.18, 1.82, 3.15, 3.5, fill=WHITE, line=BORDER)
    add_icon(slide, "pdf", 10.25, 2.16, 0.72)
    add_text(slide, "300 PDFs", 9.68, 3.15, 2.2, 0.45, size=28, bold=True, font="Poppins", align=PP_ALIGN.CENTER)
    add_text(slide, "Ready to print, email or archive", 9.62, 3.86, 2.3, 0.28, size=11, color=GRAY, align=PP_ALIGN.CENTER)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Benefits", "Practical advantages for teams handling repetitive office documents.")
    for idx, (label, icon) in enumerate(BENEFITS):
        col = idx % 3
        row = idx // 3
        x = 1.0 + col * 3.85
        y = 2.0 + row * 1.48
        add_round_rect(slide, x, y, 3.2, 1.08, fill=WHITE, line=BORDER)
        add_icon(slide, icon, x + 0.24, y + 0.26, 0.42)
        add_text(slide, label, x + 0.9, y + 0.38, 1.9, 0.25, size=14, bold=True, font="Poppins")

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Who Can Use It", "Designed for education and training organizations with high document volume.")
    for idx, (label, icon) in enumerate(AUDIENCES):
        col = idx % 3
        row = idx // 3
        x = 1.0 + col * 3.85
        y = 2.0 + row * 1.48
        add_round_rect(slide, x, y, 3.2, 1.08, fill=WHITE, line=BORDER)
        add_icon(slide, icon, x + 0.24, y + 0.26, 0.42)
        add_text(slide, label, x + 0.88, y + 0.36, 2.0, 0.3, size=13, bold=True, font="Poppins")

    # Slide 8
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)
    add_logo(slide, 0.78, 0.55, 0.9)
    add_text(slide, "Free 20-Minute Office Automation Audit", 0.92, 1.85, 6.7, 1.0, size=38, bold=True, font="Poppins")
    add_text(slide, "Book your free consultation and identify the office tasks your team can automate first.", 0.96, 3.08, 6.4, 0.52, size=17, color=GRAY)
    add_button(slide, "Book Free Consultation", 0.96, 4.05, 2.6, 0.56)
    add_text(slide, "https://aaryarushi.vercel.app", 0.98, 4.86, 3.6, 0.28, size=14, color=BLUE, bold=True)
    add_round_rect(slide, 8.58, 1.72, 2.85, 2.85, fill=WHITE, line=BORDER)
    slide.shapes.add_picture(str(ASSET_PNG / "qr_placeholder.png"), Inches(8.86), Inches(1.98), width=Inches(2.3))
    add_text(slide, "Scan to open website", 8.55, 4.86, 2.9, 0.26, size=12, color=GRAY, align=PP_ALIGN.CENTER)

    path = OUT / "AaryaRushi_Automation_Labs_Business_Presentation.pptx"
    prs.save(path)
    return path


VOICEOVER = """# Voice-over Script

Estimated duration: 90-120 seconds

## Slide 1
AaryaRushi Automation Labs helps education and training organizations save hours of manual office work every week. Our office automation solutions turn repetitive document tasks into simple, reliable workflows.

## Slide 2
Most teams still spend valuable time creating certificates, worksheets, attendance reports, ID cards, results, question papers and timetables manually. This work is repetitive, slow and easy to get wrong when the volume increases.

## Slide 3
We automate these daily workflows through focused products like AR-CERT-PRO, AR-WORKSHEET-PRO, AR-QUESTION-PRO, AR-ATTENDANCE-PRO, AR-IDCARD-PRO, AR-RESULT-PRO and AR-TIMETABLE-PRO. Each solution is designed to reduce typing, formatting and repeated checking.

## Slide 4
The process is simple. Your team keeps working with familiar Excel data. The automation engine reads the information, applies approved Word templates, generates PDFs, and delivers ready-to-use documents.

## Slide 5
For example, if an Excel file contains 300 students, your team can click Generate and let the system create all 300 PDFs automatically. The progress is visible, the outputs are organized, and the completed files are ready for printing or sharing.

## Slide 6
The result is a faster office workflow. Teams save time, reduce errors, create professional documents, complete tasks in one click, reuse templates and produce printable PDFs consistently.

## Slide 7
These solutions are useful for coaching institutes, schools, colleges, training institutes, HR teams and certification organizations that handle repeated document work every month.

## Slide 8
To get started, book a free 20-minute Office Automation Audit. We will review your current workflow and identify where automation can save the most time first. Visit https://aaryarushi.vercel.app to request your free consultation.
"""


ANIMATION_GUIDE = """# Animation Guide

Animation rules: use only Fade, Slide, Zoom and Grow. Avoid spinning, bouncing, flashy transitions and aggressive motion.

## Slide 1
- Logo: Fade, 0.4 seconds.
- Headline: Slide Up, 0.6 seconds, after logo.
- Subtitle: Fade, 0.4 seconds, after headline.
- CTA button: Grow pulse once, subtle 105% scale, 0.5 seconds.
- Dashboard visual: Fade, 0.5 seconds.

## Slide 2
- Title: Fade, 0.3 seconds.
- Problem cards: Fade one by one, 0.25 seconds each.
- Manual work illustration: Fade after the first two problem cards.

## Slide 3
- Product cards: Slide Up in sequence, 0.25 seconds each.
- Use a short 0.05 second delay between cards.

## Slide 4
- Flow items: Fade in from left to right.
- Arrows: Fade after each previous item.
- Closing statement: Fade last.

## Slide 5
- Excel card: Fade.
- Generate button: Grow once.
- Progress bar: Grow from left to right.
- Counter numbers: Fade in sequence: 25, 75, 150, 300.
- Completed status: Fade with green emphasis.
- 300 PDFs output card: Zoom in subtly.

## Slide 6
- Benefit cards: Slide Up in two rows, 0.25 seconds each.
- Keep all motion subtle and linear.

## Slide 7
- Institution cards: Fade one by one, 0.25 seconds each.
- Icons should appear with their cards, not separately.

## Slide 8
- Main offer: Slide Up.
- Consultation button: Grow pulse once.
- Website URL: Fade.
- QR placeholder: Fade in last.
"""


PRESENTER_NOTES = """# Presenter Notes

## Slide 1
Open with the core promise: fewer manual office hours and faster document output. Keep the tone practical and business-focused.

## Slide 2
Connect with the audience's real daily pain: repeated typing, formatting, checking, printing and corrections.

## Slide 3
Position each product as a focused solution, not a generic software package. Emphasize that the system is built around actual office tasks.

## Slide 4
Explain that staff can continue using familiar Excel and Word formats. This reduces training effort and makes adoption easier.

## Slide 5
Use the 300-student example as the clearest business proof. The key message is volume handled in one click.

## Slide 6
Summarize outcomes in operational language: time saved, fewer mistakes, professional outputs and reusable templates.

## Slide 7
Make the audience feel included. Education, training, HR and certification teams all face similar repetitive document workflows.

## Slide 8
Close with one action only: book the free 20-minute Office Automation Audit. Invite them to bring one repetitive workflow for review.
"""


FOLDER_STRUCTURE = """# Folder Structure

```text
business_presentation/
  AaryaRushi_Automation_Labs_Business_Presentation.pptx
  assets/
    png/
      aaryarushi_logo.png
      qr_placeholder.png
      editable icon PNG companions
    svg/
      aaryarushi_logo.svg
      editable SVG icons
  docs/
    animation_guide.md
    folder_structure.md
    presenter_notes.md
    voice_over_script.md
  source/
    create_presentation.py
```
"""


def write_docs():
    (DOCS / "voice_over_script.md").write_text(VOICEOVER, encoding="utf-8")
    (DOCS / "animation_guide.md").write_text(ANIMATION_GUIDE, encoding="utf-8")
    (DOCS / "presenter_notes.md").write_text(PRESENTER_NOTES, encoding="utf-8")
    (DOCS / "folder_structure.md").write_text(FOLDER_STRUCTURE, encoding="utf-8")


def main():
    ensure_dirs()
    save_svg_assets()
    save_png_assets()
    pptx = create_presentation()
    write_docs()
    print(f"Created {pptx}")
    print(f"Assets: {ASSET_SVG} and {ASSET_PNG}")
    print(f"Docs: {DOCS}")


if __name__ == "__main__":
    main()
